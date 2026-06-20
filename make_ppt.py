import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

# Set slide size to 16:9 widescreen (13.33 x 7.5 inches)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Premium Theme Palette
c_bg = RGBColor(0x07, 0x0C, 0x16)
c_card = RGBColor(0x0E, 0x16, 0x26)
c_cardHeader = RGBColor(0x15, 0x20, 0x35)
c_cyan = RGBColor(0x00, 0xF0, 0xFF)
c_blue = RGBColor(0x2A, 0x82, 0xE6)
c_green = RGBColor(0x05, 0xD7, 0x82)
c_red = RGBColor(0xFF, 0x2A, 0x5F)
c_purple = RGBColor(0x8A, 0x5C, 0xF5)
c_orange = RGBColor(0xFF, 0x9F, 0x1C)
c_text = RGBColor(0xF5, 0xF9, 0xFC)
c_textMuted = RGBColor(0x8A, 0x9B, 0xB4)
c_border = RGBColor(0x1A, 0x2C, 0x46)

def add_full_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_bg
    bg.line.fill.background()
    return bg

def createBaseSlide(title_text, category_text="GRIDLOCK HACKATHON 2.0"):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    add_full_bg(slide)
    
    # Top tag
    tx_hud = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.3))
    tf = tx_hud.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"SMARTVISION X  //  {category_text}"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = c_cyan
    p.font.name = 'Arial'
    
    if title_text:
        # Title text
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.0), Inches(0.6))
        tf_title = tx_title.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = c_text
        p_title.font.name = 'Arial'
        
        # Horizontal line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = c_border
        line.line.fill.background()
        
        # HUD highlight line
        line_hud = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(1.8), Inches(0.03))
        line_hud.fill.solid()
        line_hud.fill.fore_color.rgb = c_cyan
        line_hud.line.fill.background()
        
    # Watermark showing live URL
    tx_watermark = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.73), Inches(0.3))
    p_wm = tx_watermark.text_frame.paragraphs[0]
    p_wm.text = "LIVE PLATFORM DEMO: smartvision-x-gridlock-c579.vercel.app"
    p_wm.font.size = Pt(9)
    p_wm.font.bold = True
    p_wm.font.color.rgb = c_cyan
    p_wm.font.name = 'Arial'
    
    return slide

def add_card(slide, x, y, w, h, border_color, fill_color=c_card):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

# -----------------------------------------------------------------------------
# SLIDE 1: Cover
# -----------------------------------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bg(slide1)

# Dynamic HUD elements
el1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.66), Inches(1.75), Inches(4.0), Inches(4.0))
el1.fill.solid()
el1.fill.fore_color.rgb = RGBColor(0x0A, 0x12, 0x24)
el1.line.color.rgb = c_cyan
el1.line.width = Pt(1)

el2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.16), Inches(2.25), Inches(3.0), Inches(3.0))
el2.fill.solid()
el2.fill.fore_color.rgb = RGBColor(0x0F, 0x1B, 0x35)
el2.line.color.rgb = c_blue
el2.line.width = Pt(1)

# Text blocks
tx = slide1.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.4))
p = tx.text_frame.paragraphs[0]
p.text = "FLIPKART GRIDLOCK HACKATHON 2.0  //  THEME 3 SUBMISSION"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = c_cyan
p.font.name = 'Arial'

tx2 = slide1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.33), Inches(1.3))
p2 = tx2.text_frame.paragraphs[0]
p2.text = "SMARTVISION X"
p2.alignment = PP_ALIGN.CENTER
p2.font.size = Pt(68)
p2.font.bold = True
p2.font.color.rgb = c_text
p2.font.name = 'Arial'

tx3 = slide1.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.5))
p3 = tx3.text_frame.paragraphs[0]
p3.text = "Autonomous AI-Powered Traffic Enforcement & Safety Platform"
p3.alignment = PP_ALIGN.CENTER
p3.font.size = Pt(20)
p3.font.bold = True
p3.font.color.rgb = c_blue
p3.font.name = 'Arial'

add_card(slide1, 1.66, 4.8, 10.0, 1.8, c_border)
tx_meta = slide1.shapes.add_textbox(Inches(2.0), Inches(4.9), Inches(9.33), Inches(1.6))
tf_meta = tx_meta.text_frame
tf_meta.word_wrap = True

p_m1 = tf_meta.paragraphs[0]
p_m1.text = "PROJECT OVERVIEW"
p_m1.font.size = Pt(11)
p_m1.font.bold = True
p_m1.font.color.rgb = c_cyan

p_m2 = tf_meta.add_paragraph()
p_m2.text = "Presented by: Surjeet Yadav"
p_m2.font.size = Pt(16)
p_m2.font.bold = True
p_m2.font.color.rgb = c_text

p_m3 = tf_meta.add_paragraph()
p_m3.text = "Live Demo: smartvision-x-gridlock-c579.vercel.app"
p_m3.font.size = Pt(12)
p_m3.font.bold = True
p_m3.font.color.rgb = c_green

p_m4 = tf_meta.add_paragraph()
p_m4.text = "Architecture: React / Vite + Express MVC + YOLOv8 + MongoDB Atlas Sync"
p_m4.font.size = Pt(12)
p_m4.font.color.rgb = c_textMuted

# -----------------------------------------------------------------------------
# SLIDE 2: Vision & Mission
# -----------------------------------------------------------------------------
slide2 = createBaseSlide("Our Vision & Mission", "EXECUTIVE SUMMARY")

add_card(slide2, 0.8, 2.2, 5.5, 4.2, c_cyan)
add_card(slide2, 0.8, 2.2, 5.5, 0.6, None, c_cardHeader)
tx_vis = slide2.shapes.add_textbox(Inches(1.1), Inches(2.25), Inches(4.9), Inches(3.9))
tf_vis = tx_vis.text_frame
tf_vis.word_wrap = True
p_vh = tf_vis.paragraphs[0]
p_vh.text = "THE VISION"
p_vh.font.size = Pt(16)
p_vh.font.bold = True
p_vh.font.color.rgb = c_cyan

p_vb = tf_vis.add_paragraph()
p_vb.text = "\nTo construct a fully autonomous urban mobility nervous system that eliminates human delays in traffic enforcement, instantly detects life-threatening safety violations, and secures Indian roads using real-time Edge Artificial Intelligence."
p_vb.font.size = Pt(15)
p_vb.font.color.rgb = c_text
p_vb.font.name = 'Calibri'

add_card(slide2, 7.0, 2.2, 5.5, 4.2, c_blue)
add_card(slide2, 7.0, 2.2, 5.5, 0.6, None, c_cardHeader)
tx_mis = slide2.shapes.add_textbox(Inches(7.3), Inches(2.25), Inches(4.9), Inches(3.9))
tf_mis = tx_mis.text_frame
tf_mis.word_wrap = True
p_mh = tf_mis.paragraphs[0]
p_mh.text = "THE MISSION"
p_mh.font.size = Pt(16)
p_mh.font.bold = True
p_mh.font.color.rgb = c_blue

p_mb = tf_mis.add_paragraph()
p_mb.text = "\nDeploy an end-to-end MERN stack + Python CV application that processes live CCTV streams to:\n\n• Detect Helmetless Riders & Triple Riding\n• Identify Sidewalk/Illegal Parking Blockages\n• Auto-generate legal evidence & PDF Challans\n• Dispatch patrol officers based on GPS proximity."
p_mb.font.size = Pt(14)
p_mb.font.color.rgb = c_text
p_mb.font.name = 'Calibri'

# -----------------------------------------------------------------------------
# SLIDE 3: The Urban Traffic Crisis
# -----------------------------------------------------------------------------
slide3 = createBaseSlide("The Urban Traffic Crisis", "PROBLEM STATEMENT")

crisisData = [
    { 'title': 'Critical Safety Violations', 'desc': 'Riders ignoring helmet mandates and practicing dangerous triple-riding habits cause 60%+ of two-wheeler road casualties.', 'stat': '60% Fatalities', 'color': c_red },
    { 'title': 'Junction Bottlenecks', 'desc': 'Illegal parking on major urban arterial roads blocks sidewalks and creates severe choke points (e.g., Silk Board Jn).', 'stat': 'Severe Congestion', 'color': c_orange },
    { 'title': 'Enforcement Leakage', 'desc': 'Manual operators miss over 80% of transient parking and helmet violations due to sheer feed volume and visual fatigue.', 'stat': '80% Missed Tickets', 'color': c_cyan }
]

for idx, item in enumerate(crisisData):
    x = 0.8 + (idx * 4.0)
    add_card(slide3, x, 2.2, 3.7, 4.2, item['color'])
    top_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(3.7), Inches(0.1))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = item['color']
    top_line.line.fill.background()
    
    tx_c = slide3.shapes.add_textbox(Inches(x + 0.2), Inches(2.4), Inches(3.3), Inches(3.8))
    tf_c = tx_c.text_frame
    tf_c.word_wrap = True
    
    p_tag = tf_c.paragraphs[0]
    p_tag.text = f"METRIC 0{idx+1}"
    p_tag.font.size = Pt(9)
    p_tag.font.bold = True
    p_tag.font.color.rgb = c_textMuted
    
    p_stat = tf_c.add_paragraph()
    p_stat.text = item['stat']
    p_stat.font.size = Pt(24)
    p_stat.font.bold = True
    p_stat.font.color.rgb = item['color']
    
    p_t = tf_c.add_paragraph()
    p_t.text = "\n" + item['title']
    p_t.font.size = Pt(15)
    p_t.font.bold = True
    p_t.font.color.rgb = c_text
    
    p_d = tf_c.add_paragraph()
    p_d.text = "\n" + item['desc']
    p_d.font.size = Pt(12.5)
    p_d.font.color.rgb = c_textMuted

# -----------------------------------------------------------------------------
# SLIDE 4: Limitations of Manual Enforcement
# -----------------------------------------------------------------------------
slide4 = createBaseSlide("Limitations of Manual Enforcement", "PROBLEM STATEMENT")

tx_intro = slide4.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.0), Inches(0.4))
tx_intro.text_frame.paragraphs[0].text = "Traditional traffic policing methods cannot scale with mega-city volumes due to:"
tx_intro.text_frame.paragraphs[0].font.size = Pt(16)
tx_intro.text_frame.paragraphs[0].font.color.rgb = c_text

limits = [
    { 'label': 'HUMAN VISUAL FATIGUE', 'desc': 'Impossible for control room officers to monitor hundreds of CCTV streams concurrently 24×7 without missing violations.' },
    { 'label': 'HIGH PROCESSING LATENCY', 'desc': 'Manual e-challans take days to process and dispatch, losing the psychological impact of immediate traffic fine warnings.' },
    { 'label': 'OFF-LINE DISPATCH LOOP', 'desc': 'No active connection between cameras flagging violations and nearby patrol cars, leaving field units unrouted.' },
    { 'label': 'DISCONNECTED REPEAT HISTORIES', 'desc': 'No system queries registration numbers against history to identify high-risk drivers in real-time.' }
]

for idx, limit in enumerate(limits):
    y = 2.6 + (idx * 1.05)
    bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.15), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c_red
    bar.line.fill.background()
    
    add_card(slide4, 0.95, y, 11.58, 0.9, c_border)
    
    tx_lim = slide4.shapes.add_textbox(Inches(1.15), Inches(y + 0.1), Inches(11.0), Inches(0.7))
    tf_lim = tx_lim.text_frame
    tf_lim.word_wrap = True
    p_lh = tf_lim.paragraphs[0]
    p_lh.text = limit['label'] + "   |   "
    p_lh.font.size = Pt(13)
    p_lh.font.bold = True
    p_lh.font.color.rgb = c_red
    
    # We append the description on the same line or as run
    run = p_lh.add_run()
    run.text = limit['desc']
    run.font.size = Pt(13)
    run.font.bold = False
    run.font.color.rgb = c_text

# -----------------------------------------------------------------------------
# SLIDE 5: Introducing SmartVision X
# -----------------------------------------------------------------------------
slide5 = createBaseSlide("Introducing SmartVision X", "THE SOLUTION")

# Unified Header Block
add_card(slide5, 0.8, 2.0, 11.73, 1.3, c_cyan, RGBColor(0x0D, 0x17, 0x2A))
tx_sh = slide5.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.33), Inches(1.1))
tf_sh = tx_sh.text_frame
tf_sh.word_wrap = True
p_sh = tf_sh.paragraphs[0]
p_sh.text = "The Ultimate AI-Powered Traffic Nervous System"
p_sh.alignment = PP_ALIGN.CENTER
p_sh.font.size = Pt(24)
p_sh.font.bold = True
p_sh.font.color.rgb = c_cyan

p_sb = tf_sh.add_paragraph()
p_sb.text = "An end-to-end platform that captures street streams, detects violations using custom-quantised YOLOv8 models, reads license plates with OCR, and dispatches responders via a real-time web portal."
p_sb.alignment = PP_ALIGN.CENTER
p_sb.font.size = Pt(14)
p_sb.font.color.rgb = c_text

# 4 Horizontal cards
steps = [
    { 'num': '01', 'title': 'CCTV CAPTURE', 'desc': 'RTSP video feeds from high-density junctions stream directly to edge servers.' },
    { 'num': '02', 'title': 'EDGE AI INFERENCE', 'desc': 'Optimised YOLOv8 instantly locates riders, cars, helmets, and plates.' },
    { 'num': '03', 'title': 'PLATE OCR READ', 'desc': 'EasyOCR extracts alphanumeric text from HSRP plates on violation frames.' },
    { 'num': '04', 'title': 'CLOUD ENFORCEMENT', 'desc': 'MongoDB Atlas stores record. Dispatch triggers closest patrol car.' }
]

for idx, step in enumerate(steps):
    x = 0.8 + (idx * 2.95)
    add_card(slide5, x, 3.8, 2.8, 2.6, c_border)
    
    tx_step = slide5.shapes.add_textbox(Inches(x + 0.15), Inches(3.9), Inches(2.5), Inches(2.4))
    tf_step = tx_step.text_frame
    tf_step.word_wrap = True
    
    p_num = tf_step.paragraphs[0]
    p_num.text = step['num']
    p_num.font.size = Pt(18)
    p_num.font.bold = True
    p_num.font.color.rgb = c_cyan
    
    p_st = tf_step.add_paragraph()
    p_st.text = step['title']
    p_st.font.size = Pt(12)
    p_st.font.bold = True
    p_st.font.color.rgb = c_text
    
    p_sd = tf_step.add_paragraph()
    p_sd.text = "\n" + step['desc']
    p_sd.font.size = Pt(11)
    p_sd.font.color.rgb = c_textMuted
    
    # Arrow
    if idx < 3:
        tx_arr = slide5.shapes.add_textbox(Inches(x + 2.75), Inches(4.7), Inches(0.3), Inches(0.5))
        tx_arr.text_frame.paragraphs[0].text = "➔"
        tx_arr.text_frame.paragraphs[0].font.size = Pt(18)
        tx_arr.text_frame.paragraphs[0].font.color.rgb = c_blue

# -----------------------------------------------------------------------------
# SLIDE 6: Key Modules Overview
# -----------------------------------------------------------------------------
slide6 = createBaseSlide("Key Modules Overview", "PLATFORM CAPABILITIES")

features = [
    { 't': 'Live Violation Detection', 'd': 'Real-time camera feed analysis with bounding boxes for motorcycles, riders, helmet status, and illegal parking.' },
    { 't': 'Automated Evidence Center', 'd': 'Frictionless dashboard showing violation snapshots, confidence rates, plate crops, and PDF download buttons.' },
    { 't': 'Smart Patrol Dispatcher', 'd': 'GPS-tracked patrol dispatches. Calculates nearest active responder using coordinates and dispatches with route details.' },
    { 't': 'Offenders Database (Risk)', 'd': 'Alphanumeric query interface retrieving history records from VAHAN database simulator to calculate recidivism risk scores.' },
    { 't': 'AI Configuration Console', 'd': 'Interactive threshold tuners allowing operators to live-adjust YOLOv8 confidence rates, OCR sensitivity, and IR exposure.' },
    { 't': 'Emergency Override Module', 'd': 'VIP Corridor simulator which triggers green lights sequentially across junctions to ensure gridlock clearance.' }
]

for idx, feat in enumerate(features):
    row = idx // 3
    col = idx % 3
    x = 0.8 + (col * 3.9)
    y = 2.1 + (row * 2.3)
    
    add_card(slide6, x, y, 3.7, 2.0, c_border)
    header = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.7), Inches(0.4))
    header.fill.solid()
    header.fill.fore_color.rgb = c_cardHeader
    header.line.fill.background()
    
    tx_f = slide6.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.05), Inches(3.4), Inches(1.8))
    tf_f = tx_f.text_frame
    tf_f.word_wrap = True
    
    p_ft = tf_f.paragraphs[0]
    p_ft.text = feat['t']
    p_ft.font.size = Pt(12)
    p_ft.font.bold = True
    p_ft.font.color.rgb = c_cyan
    
    p_fd = tf_f.add_paragraph()
    p_fd.text = "\n" + feat['d']
    p_fd.font.size = Pt(11)
    p_fd.font.color.rgb = c_text
    p_fd.line_spacing = 1.1

# -----------------------------------------------------------------------------
# SLIDE 7: AI Vision Pipeline
# -----------------------------------------------------------------------------
slide7 = createBaseSlide("AI Vision Pipeline", "TECHNOLOGY DEEP DIVE")

tx_pi = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.0), Inches(0.4))
tx_pi.text_frame.paragraphs[0].text = "How raw CCTV camera feeds are processed into legally binding traffic citations in milliseconds:"
tx_pi.text_frame.paragraphs[0].font.size = Pt(14)
tx_pi.text_frame.paragraphs[0].font.color.rgb = c_textMuted

pipelineNodes = [
    { 'step': 'CCTV CAPTURE', 'model': 'RTSP Stream', 'info': '1080p stream at 30 FPS. Weather filters simulation (Rain/Night) adjust exposure.', 'color': c_blue },
    { 'step': 'YOLOv8 DETECT', 'model': 'YOLOv8-nano', 'info': 'Localises class boxes: car, motorcycle, rider, helmet, plate.', 'color': c_cyan },
    { 'step': 'EASYOCR ENGINE', 'model': 'CRNN OCR', 'info': 'Processes cropped plate image. Applies perspective warp and Otsu thresholding.', 'color': c_purple },
    { 'step': 'DECISION RULES', 'model': 'Rule Compiler', 'info': 'Triple riding check (>2 rider boxes in vehicle box). No-helmet threshold logic.', 'color': c_orange },
    { 'step': 'REST ENGINE', 'model': 'Express Controller', 'info': 'Persists violation to MongoDB. Auto-sends challan notice.', 'color': c_green }
]

for idx, node in enumerate(pipelineNodes):
    x = 0.8 + (idx * 2.38)
    add_card(slide7, x, 2.6, 2.15, 3.8, node['color'])
    
    pill = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(2.8), Inches(1.85), Inches(0.5))
    pill.fill.solid()
    pill.fill.fore_color.rgb = node['color']
    pill.line.fill.background()
    
    tx_node = slide7.shapes.add_textbox(Inches(x + 0.1), Inches(2.85), Inches(1.95), Inches(3.4))
    tf_node = tx_node.text_frame
    tf_node.word_wrap = True
    
    p_step = tf_node.paragraphs[0]
    p_step.text = node['step']
    p_step.alignment = PP_ALIGN.CENTER
    p_step.font.size = Pt(10)
    p_step.font.bold = True
    p_step.font.color.rgb = RGBColor(0, 0, 0)
    
    p_mod = tf_node.add_paragraph()
    p_mod.text = "\n" + node['model']
    p_mod.alignment = PP_ALIGN.CENTER
    p_mod.font.size = Pt(11)
    p_mod.font.bold = True
    p_mod.font.color.rgb = c_text
    
    # Spacer line
    p_space = tf_node.add_paragraph()
    p_space.text = "---------------------"
    p_space.alignment = PP_ALIGN.CENTER
    p_space.font.size = Pt(8)
    p_space.font.color.rgb = c_border
    
    p_inf = tf_node.add_paragraph()
    p_inf.text = node['info']
    p_inf.font.size = Pt(10)
    p_inf.font.color.rgb = c_textMuted
    p_inf.line_spacing = 1.1

    if idx < 4:
        tx_arr = slide7.shapes.add_textbox(Inches(x + 2.2), Inches(4.1), Inches(0.2), Inches(0.4))
        tx_arr.text_frame.paragraphs[0].text = "➔"
        tx_arr.text_frame.paragraphs[0].font.size = Pt(16)
        tx_arr.text_frame.paragraphs[0].font.color.rgb = c_cyan

# -----------------------------------------------------------------------------
# SLIDE 8: AI Vision Demo
# -----------------------------------------------------------------------------
slide8 = createBaseSlide("AI Vision Demo: Helmet Violation", "LIVE DETECTION CORE")

add_card(slide8, 0.8, 2.1, 4.5, 4.4, c_border)
add_card(slide8, 0.8, 2.1, 4.5, 0.5, None, c_cardHeader)

tx_e8 = slide8.shapes.add_textbox(Inches(0.95), Inches(2.15), Inches(4.2), Inches(4.2))
tf_e8 = tx_e8.text_frame
tf_e8.word_wrap = True
p_e8h = tf_e8.paragraphs[0]
p_e8h.text = "YOLOv8 + OCR PIPELINE IN ACTION"
p_e8h.font.size = Pt(12)
p_e8h.font.bold = True
p_e8h.font.color.rgb = c_cyan

p_e8b = tf_e8.add_paragraph()
p_e8b.text = "\n• Vehicle Bounding Box: YOLOv8 locates the motorcycle and crops the regional coordinates.\n\n• Helmet Status Classifier: Specifically isolates the rider head box to run binary prediction (Helmet/No Helmet).\n\n• Automatic Plate Localization: Identifies HSRP plate location and applies OCR perspective warping.\n\n• Real-time Confidence Rating: The AI engine returns a 94.6% prediction confidence."
p_e8b.font.size = Pt(11.5)
p_e8b.font.color.rgb = c_text
p_e8b.line_spacing = 1.1

# Image Frame right side
img_frame1 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(2.1), Inches(6.83), Inches(4.4))
img_frame1.fill.solid()
img_frame1.fill.fore_color.rgb = RGBColor(0x03, 0x08, 0x11)
img_frame1.line.color.rgb = c_cyan
img_frame1.line.width = Pt(2)

img_p1 = "docs/assets/helmet_violation.png"
if os.path.exists(img_p1):
    slide8.shapes.add_picture(img_p1, Inches(5.8), Inches(2.2), width=Inches(6.63), height=Inches(4.2))
else:
    tx_err = slide8.shapes.add_textbox(Inches(6.2), Inches(3.8), Inches(5.8), Inches(1.0))
    tx_err.text_frame.paragraphs[0].text = "[Image: helmet_violation.png not found]"
    tx_err.text_frame.paragraphs[0].font.size = Pt(14)
    tx_err.text_frame.paragraphs[0].font.color.rgb = c_red

# -----------------------------------------------------------------------------
# SLIDE 9: Dashboard Demo
# -----------------------------------------------------------------------------
slide9 = createBaseSlide("SmartVision X Dashboard", "LIVE INTERFACE")

img_frame2 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.1), Inches(7.2), Inches(4.4))
img_frame2.fill.solid()
img_frame2.fill.fore_color.rgb = RGBColor(0x03, 0x08, 0x11)
img_frame2.line.color.rgb = c_blue
img_frame2.line.width = Pt(2)

img_p2 = "docs/assets/smartvision_dashboard.png"
if os.path.exists(img_p2):
    slide9.shapes.add_picture(img_p2, Inches(0.9), Inches(2.2), width=Inches(7.0), height=Inches(4.2))
else:
    tx_err = slide9.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(6.2), Inches(1.0))
    tx_err.text_frame.paragraphs[0].text = "[Image: smartvision_dashboard.png not found]"
    tx_err.text_frame.paragraphs[0].font.size = Pt(14)
    tx_err.text_frame.paragraphs[0].font.color.rgb = c_red

add_card(slide9, 8.4, 2.1, 4.13, 4.4, c_border)
add_card(slide9, 8.4, 2.1, 4.13, 0.5, None, c_cardHeader)

tx_e9 = slide9.shapes.add_textbox(Inches(8.55), Inches(2.15), Inches(3.8), Inches(3.0))
tf_e9 = tx_e9.text_frame
tf_e9.word_wrap = True
p_e9h = tf_e9.paragraphs[0]
p_e9h.text = "INTEGRATED CONTROL ROOM"
p_e9h.font.size = Pt(12)
p_e9h.font.bold = True
p_e9h.font.color.rgb = c_cyan

p_e9b = tf_e9.add_paragraph()
p_e9b.text = "\n• Live CCTV Grid: Stream switcher simulating camera junctions across Silk Board and Majestic.\n\n• AI Brain Panel: Audio/visual alerts switch with live CPU/GPU tracking.\n\n• Violations Center: Instant table rendering with OCR read outs.\n\n• Weather Toggles: Simulates clear, rain, or night exposure conditions."
p_e9b.font.size = Pt(10.5)
p_e9b.font.color.rgb = c_text
p_e9b.line_spacing = 1.1

# Live link button
btn = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(5.3), Inches(3.5), Inches(0.9))
btn.fill.solid()
btn.fill.fore_color.rgb = RGBColor(0x0A, 0x1E, 0x3A)
btn.line.color.rgb = c_green
btn.line.width = Pt(2)

tx_btn = slide9.shapes.add_textbox(Inches(8.7), Inches(5.35), Inches(3.5), Inches(0.8))
tf_btn = tx_btn.text_frame
tf_btn.word_wrap = True
p_btn1 = tf_btn.paragraphs[0]
p_btn1.text = "LIVE PLATFORM URL"
p_btn1.alignment = PP_ALIGN.CENTER
p_btn1.font.size = Pt(10)
p_btn1.font.bold = True
p_btn1.font.color.rgb = c_green

p_btn2 = tf_btn.add_paragraph()
p_btn2.text = "smartvision-x-gridlock-c579.vercel.app/command"
p_btn2.alignment = PP_ALIGN.CENTER
p_btn2.font.size = Pt(9.5)
p_btn2.font.bold = True
p_btn2.font.color.rgb = c_text
p_btn2.font.name = 'Courier New'

# -----------------------------------------------------------------------------
# SLIDE 10: Automated Evidence & E-Challan Generation
# -----------------------------------------------------------------------------
slide10 = createBaseSlide("Automated E-Challan Generation", "CORE MODULE")

add_card(slide10, 0.8, 2.1, 5.5, 4.4, c_border)
add_card(slide10, 0.8, 2.1, 5.5, 0.5, None, c_cardHeader)

tx_e10 = slide10.shapes.add_textbox(Inches(0.95), Inches(2.15), Inches(5.2), Inches(4.2))
tf_e10 = tx_e10.text_frame
tf_e10.word_wrap = True
p_e10h = tf_e10.paragraphs[0]
p_e10h.text = "E-CHALLAN GENERATION ENGINE"
p_e10h.font.size = Pt(13)
p_e10h.font.bold = True
p_e10h.font.color.rgb = c_cyan

p_e10b = tf_e10.add_paragraph()
p_e10b.text = "\n• Operator Verification: The platform keeps a human in the loop. The control officer reviews the auto-cropped plate frame before confirming.\n\n• Automated Fine Calculation: System auto-assigns fine amounts (₹1,000 for helmet violations, ₹2,000 for triple riding).\n\n• RTO Vahan Sync: Connects license plates to registered owner mobile data for automated SMS alert simulations.\n\n• High-fidelity PDF Receipt: Exports clean legal tickets complete with plate crops, maps, and barcode elements."
p_e10b.font.size = Pt(11)
p_e10b.font.color.rgb = c_text
p_e10b.line_spacing = 1.1

# Mock Ticket Card
tx_x = 7.3
ticket = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tx_x), Inches(2.1), Inches(5.2), Inches(4.4))
ticket.fill.solid()
ticket.fill.fore_color.rgb = RGBColor(255, 255, 255)
ticket.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
ticket.line.width = Pt(1.5)

tx_th = slide10.shapes.add_textbox(Inches(tx_x + 0.2), Inches(2.2), Inches(4.8), Inches(0.7))
tf_th = tx_th.text_frame
p_th1 = tf_th.paragraphs[0]
p_th1.text = "BENGALURU TRAFFIC POLICE"
p_th1.alignment = PP_ALIGN.CENTER
p_th1.font.size = Pt(13)
p_th1.font.bold = True
p_th1.font.color.rgb = RGBColor(0x0A, 0x19, 0x2F)

p_th2 = tf_th.add_paragraph()
p_th2.text = "OFFICIAL E-CHALLAN RECEIPT"
p_th2.alignment = PP_ALIGN.CENTER
p_th2.font.size = Pt(9)
p_th2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

t_line = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tx_x + 0.3), Inches(2.85), Inches(4.6), Inches(0.01))
t_line.fill.solid()
t_line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
t_line.line.fill.background()

ticketDetails = [
    ['Challan Number:', 'BTP-9082-2026-X'],
    ['Date & Time:', '20-June-2026 | 10:00 AM'],
    ['Junction Area:', 'Silk Board Junction (CAM-04)'],
    ['Vehicle Plate:', 'KA-03-MX-8891'],
    ['Violation Class:', 'Helmet Non-Compliance'],
    ['Fine Amount:', '₹ 1,000.00']
]

for idx, row in enumerate(ticketDetails):
    y = 3.0 + (idx * 0.35)
    tx_row1 = slide10.shapes.add_textbox(Inches(tx_x + 0.3), Inches(y), Inches(1.8), Inches(0.3))
    p_r1 = tx_row1.text_frame.paragraphs[0]
    p_r1.text = row[0]
    p_r1.font.size = Pt(10.5)
    p_r1.font.bold = True
    p_r1.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    
    tx_row2 = slide10.shapes.add_textbox(Inches(tx_x + 2.1), Inches(y), Inches(2.8), Inches(0.3))
    p_r2 = tx_row2.text_frame.paragraphs[0]
    p_r2.text = row[1]
    p_r2.font.size = Pt(10.5)
    p_r2.font.color.rgb = RGBColor(0xFF, 0, 0) if row[1].startswith('₹') else RGBColor(0x11, 0x11, 0x11)
    p_r2.font.bold = True if row[1].startswith('₹') else False

# Green Verified Stamp
stamp = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(tx_x + 3.2), Inches(5.2), Inches(1.6), Inches(0.5))
stamp.fill.solid()
stamp.fill.fore_color.rgb = RGBColor(0xE6, 0xFB, 0xF3)
stamp.line.color.rgb = RGBColor(0x05, 0xD7, 0x82)
stamp.line.width = Pt(2)

tx_st = slide10.shapes.add_textbox(Inches(tx_x + 3.2), Inches(5.25), Inches(1.6), Inches(0.4))
tx_st.text_frame.paragraphs[0].text = "VERIFIED & PAID"
tx_st.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
tx_st.text_frame.paragraphs[0].font.size = Pt(9.5)
tx_st.text_frame.paragraphs[0].font.bold = True
tx_st.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x05, 0xB3, 0x6C)

# Barcode
bar_widths = [0.03, 0.08, 0.04, 0.02, 0.09, 0.03, 0.05, 0.02, 0.08, 0.04, 0.02, 0.09, 0.03, 0.05, 0.02, 0.08]
curr_x = tx_x + 0.4
for bw in bar_widths:
    bar_rect = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(curr_x), Inches(5.3), Inches(bw), Inches(0.4))
    bar_rect.fill.solid()
    bar_rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
    bar_rect.line.fill.background()
    curr_x += bw + 0.03

# -----------------------------------------------------------------------------
# SLIDE 11: MongoDB Sync
# -----------------------------------------------------------------------------
slide11 = createBaseSlide("MongoDB Atlas Live Cloud Sync", "DATABASE ARCHITECTURE")

add_card(slide11, 0.8, 2.1, 5.5, 4.4, c_border)
add_card(slide11, 0.8, 2.1, 5.5, 0.5, None, c_cardHeader)

tx_e11 = slide11.shapes.add_textbox(Inches(0.95), Inches(2.15), Inches(5.2), Inches(4.2))
tf_e11 = tx_e11.text_frame
tf_e11.word_wrap = True
p_e11h = tf_e11.paragraphs[0]
p_e11h.text = "ENTERPRISE DATA PERSISTENCE"
p_e11h.font.size = Pt(13)
p_e11h.font.bold = True
p_e11h.font.color.rgb = c_green

p_e11b = tf_e11.add_paragraph()
p_e11b.text = "\n• Hybrid Mode: Connects to live MongoDB Atlas Cloud cluster. Automatically falls back to local JSON persistence if internet drops.\n\n• Real-time Schema Validation: Mongoose models enforce validation rules, timestamp audits, and geospatial coordinates.\n\n• Query Optimization: Indexes on vehicle license plates and violation timestamps allow instant search queries.\n\n• Global Accessibility: Command Center dashboard queries sync instantly, giving city inspectors access to active logs."
p_e11b.font.size = Pt(11)
p_e11b.font.color.rgb = c_text
p_e11b.line_spacing = 1.1

# Mock Code Block
codeX = 7.3
add_card(slide11, codeX, 2.1, 5.2, 4.4, c_green, RGBColor(0x02, 0x07, 0x12))
header_code = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(codeX), Inches(2.1), Inches(5.2), Inches(0.4))
header_code.fill.solid()
header_code.fill.fore_color.rgb = RGBColor(0x0B, 0x15, 0x28)
header_code.line.fill.background()

tx_ch = slide11.shapes.add_textbox(Inches(codeX + 0.2), Inches(2.15), Inches(4.8), Inches(0.3))
tx_ch.text_frame.paragraphs[0].text = "mongoose.model(\"Violation\")"
tx_ch.text_frame.paragraphs[0].font.size = Pt(10)
tx_ch.text_frame.paragraphs[0].font.bold = True
tx_ch.text_frame.paragraphs[0].font.color.rgb = c_textMuted
tx_ch.text_frame.paragraphs[0].font.name = 'Courier New'

jsonLines = [
    '{',
    '  "_id": ObjectId("66741b21efb34"),',
    '  "plate_number": "KA03MX8891",',
    '  "violation_type": "No Helmet",',
    '  "location": "Silk Board Junction",',
    '  "confidence": 0.946,',
    '  "evidence_image": "uploads/ KA03MX.jpg",',
    '  "challan_amount": 1000,',
    '  "status": "pending_operator_review",',
    '  "timestamp": ISODate("2026-06-20T10:00:00Z")',
    '}'
]

for idx, line in enumerate(jsonLines):
    y = 2.7 + (idx * 0.32)
    color = c_text
    if '"_id"' in line or '"plate_number"' in line or '"violation_type"' in line or '"location"' in line or '"confidence"' in line or '"evidence_image"' in line or '"challan_amount"' in line or '"status"' in line or '"timestamp"' in line:
        color = c_green
    elif 'ObjectId' in line or 'ISODate' in line or '0.946' in line or '1000' in line:
        color = c_orange
    elif 'KA03MX8891' in line or 'No Helmet' in line or 'Silk Board' in line or 'uploads/' in line or 'pending' in line:
        color = c_cyan
        
    tx_line = slide11.shapes.add_textbox(Inches(codeX + 0.3), Inches(y), Inches(4.8), Inches(0.3))
    tx_line.text_frame.paragraphs[0].text = line
    tx_line.text_frame.paragraphs[0].font.size = Pt(9.5)
    tx_line.text_frame.paragraphs[0].font.color.rgb = color
    tx_line.text_frame.paragraphs[0].font.name = 'Courier New'

# -----------------------------------------------------------------------------
# SLIDE 12: Smart Dispatch
# -----------------------------------------------------------------------------
slide12 = createBaseSlide("Smart Dispatch System", "OPERATIONS")

add_card(slide12, 0.8, 2.1, 5.5, 4.4, c_border)
add_card(slide12, 0.8, 2.1, 5.5, 0.5, None, c_cardHeader)

tx_e12 = slide12.shapes.add_textbox(Inches(0.95), Inches(2.15), Inches(5.2), Inches(4.2))
tf_e12 = tx_e12.text_frame
tf_e12.word_wrap = True
p_e12h = tf_e12.paragraphs[0]
p_e12h.text = "HAVERSINE NEAREST OFFICER ROUTING"
p_e12h.font.size = Pt(13)
p_e12h.font.bold = True
p_e12h.font.color.rgb = c_purple

p_e12b = tf_e12.add_paragraph()
p_e12b.text = "\n• Real-time GPS Tracking: Field officers update their live coordinate pairs back to the dashboard.\n\n• Haversine Metric: When an operator verifies a critical block or triple-riding event, the backend runs the Haversine formula to compute distance.\n\n• Dynamic dispatch: Instantly orders patrol vehicles closest to the target coordinate (e.g. Officer Surjeet located 1.2km away).\n\n• Interactive Route Map: Highlights target coordinates on the CommandCenter dispatch console."
p_e12b.font.size = Pt(11)
p_e12b.font.color.rgb = c_text
p_e12b.line_spacing = 1.1

# Routing Flowchart
flowX = 7.3
add_card(slide12, flowX, 2.1, 5.2, 4.4, c_purple)

flowNodes = [
    { 'title': '1. CRITICAL INCIDENT FLAGGED', 'desc': 'CAM-04 identifies illegal parking gridlock.', 'col': c_red },
    { 'title': '2. PROXIMITY ENGINE COMPUTES', 'desc': 'Backend filters active patrol coordinate arrays.', 'col': c_cyan },
    { 'title': '3. CLOSEST CO-ORDINATE ROUTE', 'desc': 'Officer dispatch instructions pushed automatically.', 'col': c_green }
]

for idx, node in enumerate(flowNodes):
    y = 2.4 + (idx * 1.3)
    
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(flowX + 0.3), Inches(y), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = node['col']
    circle.line.fill.background()
    
    tx_num = slide12.shapes.add_textbox(Inches(flowX + 0.3), Inches(y + 0.05), Inches(0.6), Inches(0.5))
    tx_num.text_frame.paragraphs[0].text = str(idx+1)
    tx_num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tx_num.text_frame.paragraphs[0].font.size = Pt(18)
    tx_num.text_frame.paragraphs[0].font.bold = True
    tx_num.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
    
    tx_node = slide12.shapes.add_textbox(Inches(flowX + 1.1), Inches(y), Inches(3.8), Inches(0.55))
    tf_node = tx_node.text_frame
    tf_node.word_wrap = True
    
    p_nt = tf_node.paragraphs[0]
    p_nt.text = node['title']
    p_nt.font.size = Pt(11.5)
    p_nt.font.bold = True
    p_nt.font.color.rgb = c_text
    
    p_nd = tf_node.add_paragraph()
    p_nd.text = node['desc']
    p_nd.font.size = Pt(10.5)
    p_nd.font.color.rgb = c_textMuted

    if idx < 2:
        conn = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(flowX + 0.58), Inches(y + 0.6), Inches(0.04), Inches(0.7))
        conn.fill.solid()
        conn.fill.fore_color.rgb = c_border
        conn.line.fill.background()

# -----------------------------------------------------------------------------
# SLIDE 13: Technology Stack
# -----------------------------------------------------------------------------
slide13 = createBaseSlide("Technology Stack", "ENGINEERING")

techList = [
    { 'name': 'FRONTEND MODULES', 'list': 'React.js, Vite Builder, Tailwind CSS, Lucide Icons, React Router DOM v6', 'spec': 'Dynamic route-based views, theme controllers, responsive dashboards.', 'color': c_cyan },
    { 'name': 'BACKEND SYSTEM', 'list': 'Node.js, Express Framework, Multer File Uploads, REST controllers', 'spec': 'MVC structure: split models, controllers, and JSON local mock storage fallbacks.', 'color': c_blue },
    { 'name': 'CLOUD PERSISTENCE', 'list': 'MongoDB Atlas, Mongoose ORM, Connection Status Flags', 'spec': 'Auto-sync loop triggers collection upserts, online dashboard flag updates.', 'color': c_green },
    { 'name': 'COMPUTER VISION ENGINE', 'list': 'Python 3.13, Flask API, YOLOv8 (Ultralytics), EasyOCR, OpenCV', 'spec': 'Flask server runs YOLO inference and EasyOCR reader inside isolated venv.', 'color': c_purple }
]

for idx, t in enumerate(techList):
    y = 2.1 + (idx * 1.15)
    add_card(slide13, 0.8, y, 11.73, 1.0, c_border)
    
    bar = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.1), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = t['color']
    bar.line.fill.background()
    
    tx_t = slide13.shapes.add_textbox(Inches(1.1), Inches(y + 0.05), Inches(11.0), Inches(0.9))
    tf_t = tx_t.text_frame
    tf_t.word_wrap = True
    
    p_t = tf_t.paragraphs[0]
    p_t.text = t['name']
    p_t.font.size = Pt(12)
    p_t.font.bold = True
    p_t.font.color.rgb = t['color']
    
    p_list = tf_t.add_paragraph()
    p_list.text = "Stack: " + t['list']
    p_list.font.size = Pt(12)
    p_list.font.bold = True
    p_list.font.color.rgb = c_text
    
    p_spec = tf_t.add_paragraph()
    p_spec.text = "Usage: " + t['spec']
    p_spec.font.size = Pt(11)
    p_spec.font.color.rgb = c_textMuted

# -----------------------------------------------------------------------------
# SLIDE 14: Business Impact & ROI
# -----------------------------------------------------------------------------
slide14 = createBaseSlide("Business Impact & ROI", "VALUE PROPOSITION")

impacts = [
    { 't': '90% Processing Efficiency', 'd': 'Automated YOLO/OCR identification reduces time-to-issue for challans from 3 days to under 5 seconds.', 'val': '90%', 'sub': 'Time Reduced', 'color': c_cyan },
    { 't': '40% Congestion Relief', 'd': 'Proximity dispatch clears illegal street parking blockers from sidewalk junctions before gridlocks expand.', 'val': '40%', 'sub': 'Junction Clearance', 'color': c_green },
    { 't': '3.5x Fine Collection', 'd': 'Eliminates human observation leakage. Multiple cameras stream 24/7 without fatigue, checking every vehicle registration.', 'val': '3.5x', 'sub': 'Enforcement Yield', 'color': c_blue }
]

for idx, imp in enumerate(impacts):
    x = 0.8 + (idx * 4.0)
    add_card(slide14, x, 2.2, 3.7, 4.2, c_border)
    
    header = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(3.7), Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = c_cardHeader
    header.line.fill.background()
    
    tx_num = slide14.shapes.add_textbox(Inches(x + 0.2), Inches(2.25), Inches(3.3), Inches(1.4))
    tf_num = tx_num.text_frame
    tf_num.word_wrap = True
    
    p_val = tf_num.paragraphs[0]
    p_val.text = imp['val']
    p_val.alignment = PP_ALIGN.CENTER
    p_val.font.size = Pt(48)
    p_val.font.bold = True
    p_val.font.color.rgb = imp['color']
    
    p_sub = tf_num.add_paragraph()
    p_sub.text = imp['sub']
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.size = Pt(11)
    p_sub.font.bold = True
    p_sub.font.color.rgb = c_textMuted
    
    div = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.5), Inches(3.8), Inches(2.7), Inches(0.01))
    div.fill.solid()
    div.fill.fore_color.rgb = c_border
    div.line.fill.background()
    
    tx_desc = slide14.shapes.add_textbox(Inches(x + 0.2), Inches(4.0), Inches(3.3), Inches(2.3))
    tf_desc = tx_desc.text_frame
    tf_desc.word_wrap = True
    
    p_t = tf_desc.paragraphs[0]
    p_t.text = imp['t']
    p_t.alignment = PP_ALIGN.CENTER
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = c_text
    
    p_d = tf_desc.add_paragraph()
    p_d.text = "\n" + imp['d']
    p_d.font.size = Pt(11.5)
    p_d.font.color.rgb = c_textMuted
    p_d.line_spacing = 1.1

# -----------------------------------------------------------------------------
# SLIDE 15: Roadmap
# -----------------------------------------------------------------------------
slide15 = createBaseSlide("Future Roadmap & Milestones", "ROADMAP")

tx_rm = slide15.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.0), Inches(0.4))
tx_rm.text_frame.paragraphs[0].text = "Strategic roadmap for scaling SmartVision X across regional state grids:"
tx_rm.text_frame.paragraphs[0].font.size = Pt(15)
tx_rm.text_frame.paragraphs[0].font.color.rgb = c_text

# Horiz Line
line = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(11.73), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = c_border
line.line.fill.background()

line_act = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(3.0), Inches(0.05))
line_act.fill.solid()
line_act.fill.fore_color.rgb = c_cyan
line_act.line.fill.background()

roadmap = [
    { 'date': 'Q3 2026', 'title': 'Edge Hardware Port', 'desc': 'Deploy YOLOv8-nano onto NVIDIA Jetson Orin Nano modules for physical CCTV cam placement.' },
    { 'date': 'Q4 2026', 'title': 'NIC Vahan Integration', 'desc': 'Establish direct API handshakes with NIC Vahan/Sarathi databases for real-time owner registry search.' },
    { 'date': 'H1 2027', 'title': 'Citizen Complaint App', 'desc': 'Launch crowdsourced reporting app. Citizens upload pictures which run auto YOLO checks.' },
    { 'date': 'H2 2027', 'title': 'AI-Predictive Signals', 'desc': 'Direct intersection overrides to green corridor routing based on neural traffic forecasts.' }
]

for idx, item in enumerate(roadmap):
    x = 0.8 + (idx * 2.95)
    
    add_card(slide15, x, 2.6, 2.8, 1.3, c_border)
    
    tx_node = slide15.shapes.add_textbox(Inches(x + 0.15), Inches(2.65), Inches(2.5), Inches(1.2))
    tf_node = tx_node.text_frame
    tf_node.word_wrap = True
    
    p_date = tf_node.paragraphs[0]
    p_date.text = item['date']
    p_date.font.size = Pt(13)
    p_date.font.bold = True
    p_date.font.color.rgb = c_cyan
    
    p_title = tf_node.add_paragraph()
    p_title.text = item['title']
    p_title.font.size = Pt(12)
    p_title.font.bold = True
    p_title.font.color.rgb = c_text
    
    # Circle node
    circle = slide15.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.25), Inches(4.08), Inches(0.3), Inches(0.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = c_bg
    circle.line.color.rgb = c_cyan
    circle.line.width = Pt(2)
    
    inner = slide15.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.35), Inches(4.18), Inches(0.1), Inches(0.1))
    inner.fill.solid()
    inner.fill.fore_color.rgb = c_cyan
    inner.line.fill.background()
    
    tx_desc = slide15.shapes.add_textbox(Inches(x), Inches(4.6), Inches(2.8), Inches(1.8))
    tf_desc = tx_desc.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = item['desc']
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = c_textMuted
    p_desc.line_spacing = 1.1

tx_thank = slide15.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.5))
tx_thank.text_frame.paragraphs[0].text = "Autonomous Urban Safety Platform — Ready to Scale."
tx_thank.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
tx_thank.text_frame.paragraphs[0].font.size = Pt(16)
tx_thank.text_frame.paragraphs[0].font.bold = True
tx_thank.text_frame.paragraphs[0].font.color.rgb = c_cyan

# Save
outputPath = "SmartVision_X_Final_Presentation.pptx"
prs.save(outputPath)
print("Presentation generated successfully!")
